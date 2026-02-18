# -*- coding: utf-8 -*-
"""
ppt_builder.py — Robust deck builder for CrewAI multi-agent output

Key features:
- Safely merges structured dicts from Research/Analysis/Summary steps
- Extracts from JSON strings, fenced blocks, and light markdown
- Filters out raw Python/JSON dict lines so they don't appear as bullets
- Maps key_points -> insights and trends
- Deduplicates bullets, tables, and sources
"""

from __future__ import annotations
import ast
import json
import re
from typing import Any, Dict, List, Tuple, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# --------------------------------------------------------------------
# Simple regex patterns
# --------------------------------------------------------------------
BULLET = re.compile(r'^\s*[-*\u2022]\s+')  # basic bullet removal
FENCE = re.compile(r"```(?:json)?\s*(.*?)\s*```", flags=re.S | re.I)  # fenced block
URL_PATTERN = re.compile(r"(https?://[^\s)]+)")

# --------------------------------------------------------------------
# Utilities
# --------------------------------------------------------------------
def _strip(x: Any) -> str:
    if x is None:
        return ""
    return str(x).strip()

def _coerce_list(x: Any) -> List[Any]:
    if x is None:
        return []
    return list(x) if isinstance(x, list) else [x]

def _safe_json_loads(s: str):
    try:
        return json.loads(s), None
    except Exception as e:
        return None, e

def _drop_bullet(s: str) -> str:
    return BULLET.sub("", s or "").strip()

def _find_urls(s: str) -> List[str]:
    return URL_PATTERN.findall(s or "")

def _looks_like_json_dict(s: str) -> bool:
    """
    Returns True if s is a Python/JSON dict literal (e.g., "{'title': ...}")
    We use ast.literal_eval to avoid code execution.
    """
    if not isinstance(s, str):
        return False
    s = s.strip()
    if not (s.startswith("{") and s.endswith("}")):
        return False
    try:
        obj = ast.literal_eval(s)
        return isinstance(obj, dict)
    except Exception:
        return False

def _dedupe_list_keep_order(items: List[str]) -> List[str]:
    seen, out = set(), []
    for it in items:
        if not it:
            continue
        key = it.strip()
        if key not in seen:
            seen.add(key)
            out.append(it)
    return out

def _dedupe_dict_rows(rows: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    seen, out = set(), []
    for r in rows:
        # build a stable fingerprint across typical dict rows
        key = tuple(sorted((k, _strip(v)) for k, v in r.items()))
        if key not in seen:
            seen.add(key)
            out.append(r)
    return out

# --------------------------------------------------------------------
# Split helpers
# --------------------------------------------------------------------
def _split3(s: str) -> Tuple[str, str, str]:
    parts = [p.strip() for p in re.split(r"\s*[-–:;]\s*", s or "", maxsplit=2)]
    parts += ["", "", ""]
    return _strip(parts[0]), _strip(parts[1]), _strip(parts[2])

def _split_recommendation(items: List[str]) -> Tuple[List[Optional[int]], List[str], List[str]]:
    prio, act, why = [], [], []
    for it in items:
        m = re.match(r"^\[?\s*prioritet\s*(\d+)\s*\]?\s*(.+?)(?:\s*[–-]\s*(.+))?$", it, flags=re.I)
        if m:
            prio.append(int(m.group(1)))
            act.append(_strip(m.group(2)))
            why.append(_strip(m.group(3) or ""))
        else:
            segs = re.split(r"\s*[–-]\s*", it, maxsplit=1)
            act.append(_strip(segs[0]))
            why.append(_strip(segs[1] if len(segs) > 1 else ""))
            prio.append(None)
    return prio, act, why

# --------------------------------------------------------------------
# Section maps (NO + EN)
# --------------------------------------------------------------------
SECTION_MAP_NO = {
    "sammendrag": "summary",
    "trender": "trends",
    "nøkkelpunkter": "trends",
    "hovedfunn": "insights",
    "innsikt": "insights",
    "muligheter": "opportunities",
    "risiko": "risks",
    "aktører / konkurrenter": "competitors",
    "aktorer / konkurrenter": "competitors",
    "nøkkeltall": "numbers",
    "nokkelstall": "numbers",
    "anbefalinger": "recommendations",
    "kilder": "sources",
}
SECTION_MAP_EN = {
    "executive summary": "summary",
    "key trends": "trends",
    "key points": "trends",
    "highlights": "trends",
    "findings": "insights",
    "market insights": "insights",
    "opportunities": "opportunities",
    "risks": "risks",
    "competitors": "competitors",
    "competitors / actors": "competitors",
    "key numbers": "numbers",
    "recommendations": "recommendations",
    "sources": "sources",
}
SECTION_KEYS = [
    "summary",
    "trends",
    "insights",
    "opportunities",
    "risks",
    "competitors",
    "numbers",
    "recommendations",
    "sources",
]

# --------------------------------------------------------------------
# Parse "summary=..., key_points=[...]" style strings
# --------------------------------------------------------------------
def _parse_summary_kv_block(s: str) -> Optional[Dict[str, Any]]:
    if not isinstance(s, str):
        return None
    text = s.strip()
    if text.startswith("- "):
        text = text[2:].lstrip()
    if "summary=" not in text and "key_points" not in text and "recommendations" not in text:
        return None

    def _find_value(name: str) -> Optional[str]:
        m = re.search(rf'{name}\s*=\s*"(.+?)"', text, flags=re.S | re.I)
        if m:
            return m.group(1).strip()
        m = re.search(rf"{name}\s*=\s*'(.+?)'", text, flags=re.S | re.I)
        if m:
            return m.group(1).strip()
        return None

    def _find_list(name: str) -> List[str]:
        m = re.search(rf'{name}\s*=\s*\[(.+?)\]', text, flags=re.S | re.I)
        if not m:
            return []
        body = m.group(1)
        parts = [p.strip() for p in body.split(",")]
        cleaned = []
        for p in parts:
            if (p.startswith('"') and p.endswith('"')) or (p.startswith("'") and p.endswith("'")):
                p = p[1:-1].strip()
            if p:
                cleaned.append(p)
        return [x for x in cleaned if x]

    parsed = {
        "summary": _find_value("summary"),
        "key_points": _find_list("key_points"),
        "recommendations": _find_list("recommendations"),
        "sources": _find_list("sources"),
    }
    if not (parsed["summary"] or parsed["key_points"] or parsed["recommendations"] or parsed["sources"]):
        return None
    return parsed

# --------------------------------------------------------------------
# Extract outputs index
# --------------------------------------------------------------------
def _dig_outputs(result: Any):
    data = result
    if isinstance(result, dict) and isinstance(result.get("result"), dict):
        data = result["result"]

    tasks_output = []
    if isinstance(data, dict):
        tasks_output = (
            data.get("tasks_output") or
            data.get("tasks") or
            []
        )

    # Include a final string-typed answer if present
    if isinstance(data, dict):
        final_answer = (
            data.get("final_output") or
            data.get("raw") or
            data.get("text") or
            data.get("content")
        )
        if isinstance(final_answer, str):
            tasks_output = tasks_output + [{"content": final_answer}]
        elif isinstance(final_answer, dict):
            tasks_output = tasks_output + [{"result": final_answer}]

    return data, tasks_output

# --------------------------------------------------------------------
# JSON scanning in free text
# --------------------------------------------------------------------
def _brace_scan_json(text: str) -> List[Any]:
    objs, depth, start = [], 0, None
    for i, ch in enumerate(text):
        if ch == "{":
            if depth == 0:
                start = i
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0 and start is not None:
                block = text[start:i + 1]
                try:
                    objs.append(json.loads(block))
                except Exception:
                    pass
                start = None
    return objs

# --------------------------------------------------------------------
# Markdown extraction (filters dict-looking lines)
# --------------------------------------------------------------------
def _extract_from_markdown(merged: Dict[str, Any], s: str) -> None:
    lines = [ln.rstrip() for ln in (s or "").splitlines()]
    if not lines:
        return

    def norm_header(h: str) -> Optional[str]:
        x = h.strip().lower().rstrip(":")
        return SECTION_MAP_NO.get(x) or SECTION_MAP_EN.get(x)

    current = None
    buf: List[str] = []

    def filtered_cleaned(buf: List[str]) -> List[str]:
        base = [_drop_bullet(x) for x in buf]
        # remove any raw dict-looking string lines
        base = [x for x in base if not _looks_like_json_dict(x)]
        # remove empties and heavy junk
        base = [x for x in base if x]
        return base

    def flush():
        nonlocal buf, current
        if not current or not buf:
            buf = []
            return
        cleaned = filtered_cleaned(buf)
        key = current

        if key in ("trends", "insights", "opportunities", "risks", "sources"):
            merged[key].extend(cleaned)

        elif key == "competitors":
            for x in cleaned:
                a, b, c = _split3(x)
                merged["competitors"].append({"name": a, "position": b, "notes": c})

        elif key == "numbers":
            for x in cleaned:
                a, b, c = _split3(x)
                merged["numbers"].append({"metric": a, "value": b, "source": c})

        elif key == "recommendations":
            prio, act, why = _split_recommendation(cleaned)
            for i, (p, a, w) in enumerate(zip(prio, act, why), 1):
                merged["recommendations"].append({"priority": p or i, "action": a, "rationale": w})

        elif key == "summary":
            text = " ".join(cleaned).strip()
            if len(text) > len(merged["summary"]):
                merged["summary"] = text

        buf = []

    for raw in lines:
        ln = raw.strip()
        if not ln:
            continue
        # Markdown header
        if ln.startswith("#"):
            h = ln.lstrip("# ").strip()
            canon = norm_header(h)
            if canon:
                flush()
                current = canon
                buf = []
                continue
        # Plain section header (e.g., "Risks:")
        h2 = ln.lower().rstrip(":")
        canon = norm_header(h2)
        if canon:
            flush()
            current = canon
            buf = []
            continue
        # Accumulate
        if current:
            buf.append(ln)
    flush()

# --------------------------------------------------------------------
# Merge helpers
# --------------------------------------------------------------------
def _merge_research_block(merged: Dict[str, Any], research: Dict[str, Any]):
    # Trends
    tr = research.get("trends")
    if isinstance(tr, list):
        for t in tr:
            if isinstance(t, dict):
                title = _strip(t.get("title") or t.get("name"))
                why = _strip(t.get("why_it_matters") or t.get("detail") or t.get("summary"))
                ev = _strip(t.get("evidence"))
                # readable line
                line = f"{title} — {why}".strip(" —")
                if ev:
                    line += f" (evidence: {ev})"
                if line:
                    merged["trends"].append(line)

    # Numbers
    nums = research.get("numbers") or research.get("metrics")
    if isinstance(nums, list):
        for n in nums:
            if isinstance(n, dict):
                merged["numbers"].append(
                    {
                        "metric": _strip(n.get("metric") or n.get("name")),
                        "value": _strip(n.get("value") or n.get("val")),
                        "source": _strip(n.get("source") or n.get("url")),
                    }
                )

    # Competitors / actors
    comps = research.get("competitors") or research.get("actors")
    if isinstance(comps, list):
        for c in comps:
            if isinstance(c, dict):
                merged["competitors"].append(
                    {
                        "name": _strip(c.get("name")),
                        "position": _strip(c.get("position") or c.get("role") or ""),
                        "notes": _strip(c.get("notes") or c.get("summary") or ""),
                    }
                )

    # Simple lists
    for k in ("insights", "opportunities", "risks", "sources"):
        for item in _coerce_list(research.get(k)):
            s = _strip(item)
            if s:
                merged[k].append(s)

    # Recommendations
    if isinstance(research.get("recommendations"), list):
        for r in research["recommendations"]:
            if isinstance(r, dict):
                merged["recommendations"].append(
                    {
                        "priority": r.get("priority"),
                        "action": _strip(r.get("action") or r.get("what")),
                        "rationale": _strip(r.get("rationale") or r.get("why")),
                    }
                )

def _merge_json_into(merged: Dict[str, Any], obj: Dict[str, Any]) -> None:
    if not isinstance(obj, dict):
        return

    # research block first
    if isinstance(obj.get("research"), dict):
        _merge_research_block(merged, obj["research"])

    # summary
    if "summary" in obj:
        s = _strip(obj.get("summary"))
        if len(s) > len(merged["summary"]):
            merged["summary"] = s

    # ✅ Custom handling for trends (avoid stringifying dicts)
    tr = obj.get("trends")
    if isinstance(tr, dict):
        tr = [tr]  # normalize a single trend object to a list
    if isinstance(tr, list):
        for t in tr:
            if isinstance(t, dict):
                title = _strip(t.get("title") or t.get("name"))
                why   = _strip(t.get("why_it_matters") or t.get("detail") or t.get("summary"))
                ev    = _strip(t.get("evidence"))
                line  = f"{title} — {why}".strip(" —")
                if ev:
                    line += f" (evidence: {ev})"
                if line:
                    merged["trends"].append(line)
            elif isinstance(t, str):
                if not _looks_like_json_dict(t):
                    merged["trends"].append(_strip(t))

    # ✅ Keep generic merge for simple text sections only
    for key in ("insights", "opportunities", "risks", "sources"):
        for item in _coerce_list(obj.get(key)):
            s = _strip(item)
            if s:
                merged[key].append(s)

    # Map key_points (and synonyms) to insights + trends
    for k in ("key_points", "keypoints", "highlights"):
        for item in _coerce_list(obj.get(k)):
            s = _strip(item)
            if s:
                merged["insights"].append(s)
                merged["trends"].append(s)

    # Competitors
    for comp in _coerce_list(obj.get("competitors")):
        if isinstance(comp, dict):
            merged["competitors"].append(
                {
                    "name": _strip(comp.get("name")),
                    "position": _strip(comp.get("position")),
                    "notes": _strip(comp.get("notes")),
                }
            )

    # Numbers
    for n in _coerce_list(obj.get("numbers")):
        if isinstance(n, dict):
            merged["numbers"].append(
                {
                    "metric": _strip(n.get("metric")),
                    "value": _strip(n.get("value")),
                    "source": _strip(n.get("source")),
                }
            )

    # Recommendations
    for r in _coerce_list(obj.get("recommendations")):
        if isinstance(r, dict):
            merged["recommendations"].append(
                {
                    "priority": r.get("priority"),
                    "action": _strip(r.get("action")),
                    "rationale": _strip(r.get("rationale")),
                }
            )

# --------------------------------------------------------------------
# Main JSON/Markdown extractor
# --------------------------------------------------------------------
def _extract_all_json_blocks(tasks_output, also_consider=None):
    merged = {k: ([] if k != "summary" else "") for k in SECTION_KEYS}
    candidates = []

    # collect candidate strings
    for blk in tasks_output or []:
        if isinstance(blk, dict):
            for key in ("raw", "content", "text", "final_output", "result", "output", "message"):
                s = blk.get(key)
                if isinstance(s, str) and s.strip():
                    candidates.append(s)
                elif isinstance(s, (dict, list)):
                    # collect strings deeply
                    def _collect_strings_deep(obj: Any, limit: int = 20000) -> List[str]:
                        out = []
                        def walk(x):
                            if isinstance(x, str):
                                if x.strip():
                                    out.append(x.strip())
                            elif isinstance(x, dict):
                                for v in x.values():
                                    walk(v)
                            elif isinstance(x, list):
                                for v in x:
                                    walk(v)
                        walk(obj)
                        seen, res, total = set(), [], 0
                        for s2 in out:
                            if s2 not in seen:
                                seen.add(s2)
                                res.append(s2)
                                total += len(s2)
                                if total > limit:
                                    break
                        return res
                    candidates.extend(_collect_strings_deep(s))

        # artifacts (optional)
        if isinstance(blk, dict) and isinstance(blk.get("artifacts"), list):
            for a in blk["artifacts"]:
                if isinstance(a, dict):
                    for k in ("content", "text", "raw", "result"):
                        v = a.get(k)
                        if isinstance(v, str):
                            candidates.append(v)
                        elif isinstance(v, (dict, list)):
                            # same deep collection
                            def _collect_strings_deep2(obj: Any, limit: int = 20000) -> List[str]:
                                out = []
                                def walk(x):
                                    if isinstance(x, str):
                                        if x.strip():
                                            out.append(x.strip())
                                    elif isinstance(x, dict):
                                        for v in x.values():
                                            walk(v)
                                    elif isinstance(x, list):
                                        for v in x:
                                            walk(v)
                                walk(obj)
                                seen, res, total = set(), [], 0
                                for s2 in out:
                                    if s2 not in seen:
                                        seen.add(s2)
                                        res.append(s2)
                                        total += len(s2)
                                        if total > limit:
                                            break
                                return res
                            candidates.extend(_collect_strings_deep2(v))

    for extra in (also_consider or []):
        if isinstance(extra, str):
            candidates.append(extra)
        elif isinstance(extra, (dict, list)):
            # collect strings deep (reuse)
            def _collect_strings_deep3(obj: Any, limit: int = 20000) -> List[str]:
                out = []
                def walk(x):
                    if isinstance(x, str):
                        if x.strip():
                            out.append(x.strip())
                    elif isinstance(x, dict):
                        for v in x.values():
                            walk(v)
                    elif isinstance(x, list):
                        for v in x:
                            walk(v)
                walk(obj)
                seen, res, total = set(), [], 0
                for s2 in out:
                    if s2 not in seen:
                        seen.add(s2)
                        res.append(s2)
                        total += len(s2)
                        if total > limit:
                            break
                return res
            candidates.extend(_collect_strings_deep3(extra))

    # Fenced JSON blocks
    for s in list(candidates):
        for block in FENCE.findall(s):
            cleaned = block.strip()
            # scan brace blocks
            for obj in _brace_scan_json(cleaned):
                _merge_json_into(merged, obj)
            obj, _ = _safe_json_loads(cleaned)
            if obj:
                _merge_json_into(merged, obj)

    # Bare JSON bodies (or embedded JSON objects)
    for s in list(candidates):
        ss = s.strip()
        if (ss.startswith("{") and ss.endswith("}")) or (ss.startswith("[") and ss.endswith("]")):
            obj, _ = _safe_json_loads(ss)
            if obj:
                _merge_json_into(merged, obj)
            continue
        for obj in _brace_scan_json(ss):
            _merge_json_into(merged, obj)

    # Parse "summary=..., key_points=[...]" style
    for s in list(candidates):
        kv = _parse_summary_kv_block(s)
        if kv:
            if kv.get("summary") and len(kv["summary"]) > len(merged["summary"]):
                merged["summary"] = kv["summary"]
            for p in kv.get("key_points", []):
                if p:
                    merged["insights"].append(p)
                    merged["trends"].append(p)
            for rec in kv.get("recommendations", []):
                if rec:
                    merged["recommendations"].append({"priority": None, "action": rec, "rationale": ""})
            for u in kv.get("sources", []):
                if u:
                    merged["sources"].append(u)

    # Markdown fallback (with dict-line filter)
    for s in candidates:
        _extract_from_markdown(merged, s)

    # URLs from all candidates → sources
    for u in _find_urls("\n".join(candidates)):
        merged["sources"].append(u)

    return merged

# --------------------------------------------------------------------
# Slide builders
# --------------------------------------------------------------------
def _add_bullet_slide(prs, title, bullets, size=18):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()

    # Deduplicate & filter empties
    bullets = _dedupe_list_keep_order([_strip(b) for b in bullets or []])

    if not bullets:
        tf.text = "No data available."
        return

    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(line)
        p.level = 0
        p.font.size = Pt(size)

def _add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = title

    # Deduplicate rows
    rows = _dedupe_dict_rows(rows or [])

    left, top, width, height = Inches(0.6), Inches(1.6), Inches(9.0), Inches(1.0)
    n_rows = max(2, 1 + len(rows))
    n_cols = len(headers)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    # headers
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.LEFT

    # rows
    if rows:
        for i, r in enumerate(rows[: n_rows - 1], 1):
            for j in range(n_cols):
                table.cell(i, j).text = str(r.get(headers[j].lower(), r.get(headers[j], "")))
    else:
        table.cell(1, 0).text = "No structured data available."

# --------------------------------------------------------------------
# Top-level API
# --------------------------------------------------------------------
def create_multislide_pptx(result: Dict[str, Any], topic: str, file_path: str):
    """
    Returns file_path on success; returns {"error": "..."} on failure.
    """
    try:
        # -------------------------------------------------------
        # Extract data
        # -------------------------------------------------------
        data, tasks_output = _dig_outputs(result)
        also_consider = []
        if isinstance(data, dict):
            for k in ("summary", "final_output", "raw", "content", "text"):
                v = data.get(k)
                if isinstance(v, str) and v.strip():
                    also_consider.append(v)
            # scan for more candidate strings
            def _collect_strings_deep(obj: Any, limit: int = 20000) -> List[str]:
                out = []
                def walk(x):
                    if isinstance(x, str):
                        if x.strip():
                            out.append(x.strip())
                    elif isinstance(x, dict):
                        for v in x.values():
                            walk(v)
                    elif isinstance(x, list):
                        for v in x:
                            walk(v)
                walk(obj)
                seen, res, total = set(), [], 0
                for s2 in out:
                    if s2 not in seen:
                        seen.add(s2)
                        res.append(s2)
                        total += len(s2)
                        if total > limit:
                            break
                return res
            also_consider.extend(_collect_strings_deep(data)[:50])

        sections = _extract_all_json_blocks(tasks_output, also_consider)

        # Merge structured dicts directly
        if isinstance(data, dict):
            _merge_json_into(sections, data)
        for blk in (tasks_output or []):
            if isinstance(blk, dict):
                _merge_json_into(sections, blk)
                inner = blk.get("result")
                if isinstance(inner, dict):
                    _merge_json_into(sections, inner)

        # Ensure summary fallback
        if not sections["summary"]:
            sections["summary"] = _strip(data.get("summary")) if isinstance(data, dict) else ""
            if not sections["summary"]:
                sections["summary"] = "No summary available."

        # Repair/derive topic if caller passed a bogus one
        if not topic or topic.strip().lower() == "string":
            topic_guess = ""
            if isinstance(result, dict):
                topic_guess = _strip(result.get("topic", "")) or topic_guess
            if isinstance(data, dict) and not topic_guess:
                topic_guess = _strip(data.get("topic", ""))
            topic = topic_guess or topic or "Market Insights Report"

        # Deduplicate final sections
        for k in ("trends", "insights", "opportunities", "risks", "sources"):
            sections[k] = _dedupe_list_keep_order([_strip(x) for x in sections[k]])

        sections["competitors"] = _dedupe_dict_rows(sections["competitors"])
        sections["numbers"]     = _dedupe_dict_rows(sections["numbers"])

        # Recommendations as list of strings
        rec_lines = []
        for r in sections["recommendations"]:
            if not isinstance(r, dict):
                # accept "action" only if passed as string
                rec_lines.append(_strip(str(r)))
                continue
            pr = r.get("priority")
            prefix = f"{pr}) " if isinstance(pr, int) else ""
            action = _strip(r.get("action"))
            rationale = _strip(r.get("rationale"))
            rec_lines.append(f"{prefix}{action} — Why: {rationale}".strip(" —"))

        rec_lines = _dedupe_list_keep_order(rec_lines)

        # -------------------------------------------------------
        # Create PPT
        # -------------------------------------------------------
        prs = Presentation()

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Multi‑Agent Insights Report"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = topic

        # Executive Summary
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Executive Summary"
        tf = slide.placeholders[1].text_frame
        tf.clear()
        tf.paragraphs[0].text = sections["summary"]

        # Other slides
        _add_bullet_slide(prs, "Key Trends", sections["trends"])
        _add_bullet_slide(prs, "Market Insights", sections["insights"])
        _add_bullet_slide(prs, "Opportunities", sections["opportunities"])
        _add_bullet_slide(prs, "Risks", sections["risks"])

        comp_rows = [
            {"name": _strip(c.get("name")), "position": _strip(c.get("position")), "notes": _strip(c.get("notes"))}
            for c in sections["competitors"]
        ]
        _add_table_slide(prs, "Competitors / Actors", ["Name", "Position", "Notes"], comp_rows)

        num_rows = [
            {"metric": _strip(n.get("metric")), "value": _strip(n.get("value")), "source": _strip(n.get("source"))}
            for n in sections["numbers"]
        ]
        _add_table_slide(prs, "Key Numbers", ["Metric", "Value", "Source"], num_rows)

        _add_bullet_slide(prs, "Recommendations", rec_lines)
        _add_bullet_slide(prs, "Sources", sections["sources"], size=16)

        # -------------------------------------------------------
        # Save
        # -------------------------------------------------------
        try:
            prs.save(file_path)
        except Exception as e:
            import traceback
            print("ERROR: Failed to save PPT file:", repr(e))
            traceback.print_exc()
            return {
                "error": "ppt_build_failed",
                "details": f"Could not save pptx: {type(e).__name__}: {e}"
            }
        return file_path

    except Exception as e:
        import traceback
        print("ERROR inside create_multislide_pptx:", repr(e))
        traceback.print_exc()
        return {
            "error": "ppt_build_failed",
            "details": f"{type(e).__name__}: {e}"
        }
